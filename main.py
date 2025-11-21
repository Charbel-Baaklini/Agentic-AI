import os
import io
import json
import urllib.parse
from typing import List, Dict, Any, Optional

from dotenv import load_dotenv
from openai import OpenAI
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

# =========================
# ENV + CLIENT SETUP
# =========================

load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
SERPAPI_API_KEY = os.getenv("SERPAPI_API_KEY")

# Mock mode flag: if any required key is missing, we switch to mock mode
MOCK_MODE: bool = False

if not OPENAI_API_KEY or not SERPAPI_API_KEY:
    MOCK_MODE = True
    print(
        "[Config] OPENAI_API_KEY or SERPAPI_API_KEY is missing. "
        "Running in MOCK MODE (no external API calls)."
    )
    client = None
else:
    client = OpenAI()

OUTPUT_DIR = "outputs"
IMAGES_DIR = os.path.join(OUTPUT_DIR, "images")
os.makedirs(IMAGES_DIR, exist_ok=True)


# =========================
# HELPERS
# =========================

def clean_json_string(s: str) -> str:
    """
    Remove markdown fences if the model wraps JSON in ```json ... ``` etc.
    Also strips optional leading language tags like 'json'.
    """
    s = s.strip()
    if s.startswith("```"):
        # Remove leading ``` block
        parts = s.split("```", 1)
        if len(parts) > 1:
            s = parts[1].strip()
            # Remove possible trailing ```
            if "```" in s:
                s = s.split("```", 1)[0].strip()

    # If there's a leading word like 'json' before the '{' or '[',
    # remove that first line.
    stripped = s.lstrip()
    if not (stripped.startswith("{") or stripped.startswith("[")):
        lines = s.splitlines()
        if len(lines) > 1:
            s = "\n".join(lines[1:]).strip()

    return s


def safe_json_parse(s: str, fallback: Any) -> Any:
    """
    Try to parse JSON and return fallback on any error.
    """
    try:
        s_clean = clean_json_string(s)
        return json.loads(s_clean)
    except Exception:
        return fallback


# =========================
# OPENAI FUNCTIONS (with mock fallbacks)
# =========================

def extract_main_topic(prompt: str) -> str:
    """
    Extracts a single main topic (1–2 words) from a long prompt
    to be used for image search and styling.
    """
    if MOCK_MODE:
        print("[Mock] extract_main_topic -> 'Demo Topic'")
        return "Demo Topic"

    system_msg = (
        "You will be given a long presentation prompt.\n"
        "Return ONLY the main concrete topic as 1 or at most 2 words.\n"
        "No punctuation, no explanations."
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": prompt},
        ],
        max_tokens=8,
        temperature=0.0,
    )
    topic = resp.choices[0].message.content.strip()
    topic = topic.replace('"', "").replace("'", "").strip()
    # If somehow empty, fall back to a generic word
    return topic or "topic"


def generate_slide_structure(prompt: str) -> Dict[str, Any]:
    """
    Generates a slide structure as JSON using chat completions.
    Schema:
    {
      "slides": [
        {
          "title": "Short Clean Title",
          "bullet_points": [
            "First bullet",
            "Second bullet"
          ]
        },
        ...
      ]
    }
    """
    if MOCK_MODE:
        print("[Mock] generate_slide_structure -> simple canned slides")
        return {
            "slides": [
                {
                    "title": "Demo Overview",
                    "bullet_points": [
                        "This is mock content because API keys are missing.",
                        f"Original prompt: {prompt[:80]}...",
                        "Replace MOCK_MODE by setting valid API keys in .env.",
                    ],
                },
                {
                    "title": "How It Works (Mock)",
                    "bullet_points": [
                        "No calls to OpenAI or SerpAPI.",
                        "Slides are generated locally with fixed text.",
                        "Useful for testing the pipeline without network access.",
                    ],
                },
            ]
        }

    system_msg = (
        "You are a presentation designer. Create a JSON object with key 'slides'. "
        "'slides' is a list of slides. Each slide has:\n"
        "- 'title': a very short, clean title (2–6 words).\n"
        "- 'bullet_points': 3–6 concise bullet points.\n"
        "Return ONLY valid JSON, no explanations, no extra text.\n"
        "Example:\n"
        "{\n"
        '  \"slides\": [\n'
        "    {\n"
        '      \"title\": \"Water Basics\",\n'
        '      \"bullet_points\": [\n'
        '        \"Water covers 71% of Earth\",\n'
        '        \"Essential for all known life\"\n'
        "      ]\n"
        "    }\n"
        "  ]\n"
        "}"
    )

    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": prompt},
        ],
        max_tokens=800,
        temperature=0.5,
    )

    raw_content = resp.choices[0].message.content
    fallback = {
        "slides": [
            {
                "title": "Overview",
                "bullet_points": [prompt],
            }
        ]
    }
    data = safe_json_parse(raw_content, fallback)

    # Ensure correct structure
    if not isinstance(data, dict) or "slides" not in data:
        return fallback

    slides = data.get("slides", [])
    norm_slides = []
    for s in slides:
        title = str(s.get("title", "Overview")).strip()
        bullets = s.get("bullet_points", [])
        if not isinstance(bullets, list):
            bullets = [str(bullets)]
        # make sure bullets are strings
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
        if not bullets:
            bullets = ["(content)"]
        norm_slides.append({"title": title, "bullet_points": bullets})

    if not norm_slides:
        norm_slides = fallback["slides"]

    return {"slides": norm_slides}


def choose_topic_color(topic: str) -> RGBColor:
    """
    Ask the LLM to choose a color related to the topic.
    Returns RGBColor parsed from hex or rgb(); if parsing fails,
    falls back to a neutral blue.
    """
    if MOCK_MODE:
        print("[Mock] choose_topic_color -> neutral blue")
        return RGBColor(0, 102, 204)

    system_msg = (
        "You will be given a topic.\n"
        "Return ONLY one color that best represents that topic.\n"
        "The color must be given in one of these formats:\n"
        "- Hex format like '#FFAA00'\n"
        "- rgb(r,g,b)\n"
        "No other text. No explanation."
    )

    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": topic},
        ],
        max_tokens=20,
        temperature=0.7,
    )

    color_str = resp.choices[0].message.content.strip().lower()

    # Try hex format: #rrggbb
    if color_str.startswith("#") and len(color_str) == 7:
        try:
            r = int(color_str[1:3], 16)
            g = int(color_str[3:5], 16)
            b = int(color_str[5:7], 16)
            return RGBColor(r, g, b)
        except Exception:
            pass

    # Try rgb(r,g,b)
    if color_str.startswith("rgb"):
        try:
            inner = color_str.replace("rgb", "").replace("(", "").replace(")", "")
            parts = inner.split(",")
            r, g, b = [int(p.strip()) for p in parts]
            return RGBColor(r, g, b)
        except Exception:
            pass

    # Fallback neutral color if the model misbehaves
    print(f"[Style] Could not parse color '{color_str}', using fallback.")
    return RGBColor(0, 102, 204)


# =========================
# MINI AGENT: PLAN → ACT → VERIFY
# =========================

def agent_plan_presentation(user_prompt: str) -> Dict[str, Any]:
    """
    PLAN step: Ask the model for a high-level plan of the presentation.
    Returns a JSON like:
    {
      "goal": "...",
      "sections": ["Intro", "Part 1", "Part 2", "Conclusion"],
      "estimated_slide_count": 4
    }
    """
    if MOCK_MODE:
        print("[Mock] agent_plan_presentation -> simple mock plan")
        return {
            "goal": f"Explain the main ideas of: {user_prompt[:80]}...",
            "sections": ["Overview", "Details", "Conclusion"],
            "estimated_slide_count": 3,
        }

    system_msg = (
        "You are an assistant planning a slide presentation.\n"
        "Given the user's prompt, return a JSON object with:\n"
        "- 'goal': one sentence describing the main goal of the presentation\n"
        "- 'sections': a list of short section names\n"
        "- 'estimated_slide_count': an integer between 1 and 10\n"
        "Return ONLY valid JSON."
    )

    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_prompt},
        ],
        max_tokens=200,
        temperature=0.0,
    )

    raw = resp.choices[0].message.content
    fallback = {
        "goal": user_prompt[:120],
        "sections": ["Overview"],
        "estimated_slide_count": 3,
    }
    plan = safe_json_parse(raw, fallback)
    print(f"[Agent][Plan] Goal: {plan.get('goal')}")
    print(f"[Agent][Plan] Sections: {plan.get('sections')}")
    print(f"[Agent][Plan] Estimated slides: {plan.get('estimated_slide_count')}")
    return plan


def agent_verify_and_refine_slides(
    user_prompt: str,
    plan: Dict[str, Any],
    slide_data: Dict[str, Any],
) -> Dict[str, Any]:
    """
    VERIFY step: Ask the model to check if the generated slides
    match the user prompt and the initial plan.
    It can optionally tweak titles or add one missing slide.
    """
    if MOCK_MODE:
        print("[Mock] agent_verify_and_refine_slides -> returns slides unchanged")
        return slide_data

    system_msg = (
        "You are verifying a set of slides against a user prompt and a plan.\n"
        "You will be given:\n"
        "- the original user_prompt\n"
        "- a 'plan' JSON with goal/sections/estimated_slide_count\n"
        "- the current 'slides' JSON (titles + bullet_points)\n\n"
        "Task:\n"
        "1. Check if the slides cover the main goal and the key sections.\n"
        "2. If coverage is good, return the slides unchanged.\n"
        "3. If something important is missing, you may:\n"
        "   - Improve some slide titles slightly, or\n"
        "   - Add at most ONE extra slide to cover a missing key point.\n\n"
        "Return ONLY JSON with the same schema as 'slides' from the input:\n"
        "{ \"slides\": [ { \"title\": ..., \"bullet_points\": [...] }, ... ] }"
    )

    verify_input = {
        "user_prompt": user_prompt,
        "plan": plan,
        "slides": slide_data.get("slides", []),
    }

    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_msg},
            {
                "role": "user",
                "content": json.dumps(verify_input),
            },
        ],
        max_tokens=800,
        temperature=0.2,
    )

    raw = resp.choices[0].message.content
    refined = safe_json_parse(raw, slide_data)

    # Basic sanity check
    if not isinstance(refined, dict) or "slides" not in refined:
        print("[Agent][Verify] Using original slides (refine failed).")
        return slide_data

    print(f"[Agent][Verify] Slide count after verification: {len(refined.get('slides', []))}")
    return refined


# =========================
# IMAGE SEARCH (SerpAPI only, single attempt, mocked if needed)
# =========================

def search_image_for_topic(topic: str) -> Optional[str]:
    """
    Use SerpAPI to search ONLY Google Images for the given topic.
    Single attempt (no retries).
    Returns the URL of the first image result, or None.
    """
    if MOCK_MODE:
        print("[Mock] search_image_for_topic -> None (no image in mock mode)")
        return None

    # Make the query more "safe" / generic to help SerpAPI/Google if needed
    query = f"{topic}"
    print(f"[Image] Searching Google Images for query: {query!r}")
    encoded_query = urllib.parse.quote(query)

    url = (
        "https://serpapi.com/search.json"
        f"?engine=google_images&q={encoded_query}&api_key={SERPAPI_API_KEY}"
    )

    try:
        resp = requests.get(
            url,
            timeout=20,
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
        )
        resp.raise_for_status()
        data = resp.json()
    except Exception as e:
        print(f"[Image] Error fetching image search results: {e}")
        return None

    results = data.get("images_results", [])
    print(f"[Image] Found {len(results)} image results.")
    if not results:
        print("[Image] No image results found.")
        return None

    # Find first result that has some usable URL
    for idx, r in enumerate(results):
        image_url = r.get("original") or r.get("thumbnail") or r.get("image")
        if image_url:
            print(f"[Image] Selected image #{idx} URL: {image_url}")
            return image_url

    print("[Image] No suitable image URL found in results.")
    return None


def download_image(image_url: str, filename: str = "slide_image.png") -> Optional[str]:
    """
    Downloads the image from the given URL and saves it under IMAGES_DIR.
    Returns the local file path, or None if something fails.
    """
    if not image_url:
        print("[Image] No image URL provided, skipping download.")
        return None

    try:
        resp = requests.get(
            image_url,
            timeout=30,
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
        )
        resp.raise_for_status()
        img = Image.open(io.BytesIO(resp.content))
    except Exception as e:
        print(f"[Image] Error downloading image: {e}")
        return None

    os.makedirs(OUTPUT_DIR, exist_ok=True)

    out_path = os.path.join(IMAGES_DIR, filename)
    try:
        img.convert("RGB").save(out_path, "PNG")
        print(f"[Image] Saved image to: {out_path}")
        return out_path
    except Exception as e:
        print(f"[Image] Error saving image: {e}")
        return None


# =========================
# POWERPOINT GENERATION
# =========================

def build_presentation(
    slide_data: Dict[str, Any],
    image_path: Optional[str],
    title_color: Optional[RGBColor] = None
) -> str:
    """
    Builds a PowerPoint presentation from slide_data, plus
    an extra slide containing only the image if provided.
    """
    prs = Presentation()

    slides: List[Dict[str, Any]] = slide_data.get("slides", [])
    if not slides:
        slides = [{"title": "Overview", "bullet_points": ["(no content)"]}]

    # Content slides
    for s in slides:
        title_text = s.get("title", "Overview")
        bullets = s.get("bullet_points", [])

        slide_layout = prs.slide_layouts[1]  # Title + content
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        # ----- title text + style -----
        title_shape.text = title_text
        title_tf = title_shape.text_frame
        if title_tf.paragraphs and title_tf.paragraphs[0].runs:
            title_run = title_tf.paragraphs[0].runs[0]
        else:
            # Fallback: add a run if missing
            title_p = title_tf.paragraphs[0]
            title_run = title_p.add_run()
            title_run.text = title_text

        title_run.font.bold = True
        title_run.font.size = Pt(40)
        if title_color is not None:
            title_run.font.color.rgb = title_color
        else:
            title_run.font.color.rgb = RGBColor(0, 102, 204)
        # -------------------------------

        tf = body_shape.text_frame
        tf.clear()

        if bullets:
            tf.text = str(bullets[0])
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = str(bullet)
                p.level = 0

    # Separate image-only slide
    if image_path and os.path.exists(image_path):
        blank_layout = prs.slide_layouts[6]  # Blank
        img_slide = prs.slides.add_slide(blank_layout)

        # Center-ish image with fixed size
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        img_slide.shapes.add_picture(image_path, left, top, width=width)
        print("[PPTX] Added separate image slide.")

    # Save presentation
    output_path = os.path.join(OUTPUT_DIR, "output_deck.pptx")
    prs.save(output_path)
    print(f"[PPTX] Presentation saved to: {output_path}")
    return output_path


# =========================
# ORCHESTRATION
# =========================

def generate_presentation(user_prompt: str) -> str:
    """
    High-level orchestration function (with a mini agent loop):
    PLAN:
      - Plan high-level structure from the user prompt.
    ACT:
      - Generate slide structure via OpenAI or mock.
      - Extract main topic.
      - Search & download image (SerpAPI only, single attempt; mocked to none).
    VERIFY:
      - Verify and optionally refine the slide structure (or leave as-is in mock).
    Then build PPTX with content slides + separate image slide.
    """
    # PLAN
    print("[Agent] Planning presentation...")
    plan = agent_plan_presentation(user_prompt)

    # ACT - content
    print("\n[1/3] Generating slide content...")
    slide_content = generate_slide_structure(user_prompt)

    # VERIFY
    print("[Agent] Verifying and refining slides...")
    slide_content = agent_verify_and_refine_slides(user_prompt, plan, slide_content)

    # ACT - topic + style + image
    print("[2/3] Extracting main topic for image search and styling...")
    topic = extract_main_topic(user_prompt)
    print(f"    Main topic: {topic}")

    print("[Style] Selecting a color...")
    title_color = choose_topic_color(topic)

    print("[2/3] Searching and downloading image (if not in mock mode)...")
    image_url = search_image_for_topic(topic)
    image_path = download_image(image_url) if image_url else None
    if not image_path:
        print("[Image] No image will be added.")

    # BUILD
    print("[3/3] Building PowerPoint...")
    pptx_path = build_presentation(slide_content, image_path, title_color=title_color)

    return pptx_path


# =========================
# MAIN
# =========================

def main():
    print(">> ")
    user_prompt = input("Enter the prompt for your PowerPoint slide: ").strip()
    if not user_prompt:
        print("No prompt provided, exiting.")
        return

    pptx_path = generate_presentation(user_prompt)
    print(f"\nDone! Generated presentation: {pptx_path}")


if __name__ == "__main__":
    main()
